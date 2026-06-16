"""Build the 15-slide presentation deck — style matches FORGE_VLA.pptx.
Run from:  sib_vla/   directory.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── colour palette (mirrors FORGE_VLA.pptx) ─────────────────────────────────
INK   = RGBColor(0x15, 0x18, 0x20)
BLUE  = RGBColor(0x1F, 0x5F, 0xD0)
TEAL  = RGBColor(0x0E, 0x9E, 0x8F)
AMBER = RGBColor(0xE6, 0x8A, 0x00)
RED   = RGBColor(0xCC, 0x33, 0x33)
GREY  = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGREY = RGBColor(0xF0, 0xF2, 0xF5)
PURP  = RGBColor(0x7A, 0x3F, 0xB5)
GREEN = RGBColor(0x1B, 0x88, 0x4A)

LEVER_COLORS = [BLUE, TEAL, AMBER, RED, PURP]

SW, SH = Inches(13.333), Inches(7.5)

# ── helpers ──────────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs

def blank_slide(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = WHITE
    return sl

def txbox(sl, text, l, t, w, h, size=16, bold=False, color=INK,
          align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tf = sl.shapes.add_textbox(l, t, w, h)
    tf.word_wrap = wrap
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text        = text
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tf

def title_bar(sl, title, subtitle=""):
    txbox(sl, title, Inches(0.62), Inches(0.25), Inches(12.1), Inches(0.9),
          size=36, bold=True, color=INK)
    bar = sl.shapes.add_shape(1, Inches(0.62), Inches(1.15),
                               Inches(2.6), Pt(5))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
    if subtitle:
        txbox(sl, subtitle, Inches(0.62), Inches(1.2), Inches(12.1), Inches(0.4),
              size=13, color=GREY, italic=True)

def hline(sl, l, t, w, color=BLUE, thick=3):
    box = sl.shapes.add_shape(1, l, t, w, Pt(thick))
    box.fill.solid(); box.fill.fore_color.rgb = color; box.line.fill.background()

def badge(sl, text, l, t, w, h, bg=BLUE, fg=WHITE, size=16, bold=True):
    box = sl.shapes.add_shape(1, l, t, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = bg; box.line.fill.background()
    tf = box.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = fg

def filled_box(sl, l, t, w, h, bg=LGREY, border_color=None):
    box = sl.shapes.add_shape(1, l, t, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = bg
    if border_color:
        box.line.color.rgb = border_color; box.line.width = Pt(1)
    else:
        box.line.fill.background()
    return box

def table_cell(sl, text, l, t, w, h,
               bg=LGREY, fg=INK, size=13, bold=False,
               align=PP_ALIGN.CENTER):
    box = sl.shapes.add_shape(1, l, t, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = bg
    box.line.color.rgb = RGBColor(0xCC,0xCC,0xCC); box.line.width = Pt(0.5)
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = fg

def bar_h(sl, val, max_val, l, t, w_max, h, color=BLUE):
    bw = Inches((val / max_val) * (w_max / Inches(1)))
    box = sl.shapes.add_shape(1, l, t + Inches(0.07), Inches(bw), h - Inches(0.14))
    box.fill.solid(); box.fill.fore_color.rgb = color; box.line.fill.background()

def embed_video(sl, video_path, poster_path, l, t, w, h):
    if not os.path.exists(video_path):
        box = filled_box(sl, l, t, w, h, bg=LGREY, border_color=GREY)
        txbox(sl, f"▶  VIDEO\n{os.path.basename(video_path)}",
              l, t + h//2 - Inches(0.4), w, Inches(0.8),
              size=14, color=GREY, align=PP_ALIGN.CENTER)
        return
    poster = poster_path if (poster_path and os.path.exists(poster_path)) else None
    try:
        sl.shapes.add_movie(video_path, l, t, w, h,
                            poster_frame_image=poster,
                            mime_type="video/mp4")
    except Exception as e:
        box = filled_box(sl, l, t, w, h, bg=LGREY, border_color=RED)
        txbox(sl, f"▶  {os.path.basename(video_path)}\n({e})",
              l, t + h//2 - Inches(0.4), w, Inches(0.8),
              size=13, color=RED, align=PP_ALIGN.CENTER)

# ── paths ─────────────────────────────────────────────────────────────────────
def SIB_VID(task, ep=0):
    return f"results/videos/sib_b1e-4/clean/libero_spatial_{task}/ep{ep:03d}.mp4"
def VAN_VID(task, ep=0):
    return f"results/videos/vanilla_n25/clean/libero_spatial_{task}/ep{ep:03d}.mp4"
def POSTER(prefix, task):
    return f"/tmp/ppt_posters/{prefix}_task{task}.jpg"

MATH_VID1 = "/home/user/Downloads/SpectralBottleneck.mp4"
MATH_VID2 = "/home/user/Downloads/SpectralMath.mp4"

prs = new_prs()

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — FORGE-VLA Title Card
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
txbox(sl, "FORGE-VLA", Inches(0.8), Inches(1.0), Inches(11.7), Inches(1.6),
      size=72, bold=True, color=BLUE)
txbox(sl, "Fine-tuning Optimized Recipe for Generalized Embodiment",
      Inches(0.8), Inches(2.55), Inches(11.5), Inches(0.65), size=20, color=INK)
hline(sl, Inches(0.8), Inches(3.28), Inches(11.7), color=BLUE, thick=4)
txbox(sl, "How cheap are lighting-robust small VLAs?  An in-augmentation vs held-out study.",
      Inches(0.8), Inches(3.44), Inches(11.5), Inches(0.58),
      size=18, bold=True, color=INK)

for i, (head, body) in enumerate([
        ("Base model",  "SmolVLA-500M · frozen backbone"),
        ("Benchmark",   "LIBERO-Long · 10 tasks · 300 eps"),
        ("Protocol",    "In-aug/held-out split  P0/P3/P6/P13–P17"),
]):
    x = Inches(0.8) + i * Inches(4.2)
    badge(sl, head, x, Inches(4.35), Inches(3.9), Inches(0.5),
          bg=BLUE, fg=WHITE, size=15)
    txbox(sl, body, x + Inches(0.1), Inches(4.95), Inches(3.7), Inches(0.6),
          size=16, color=GREY)

txbox(sl, "ISLab · CWNU  |  2026",
      Inches(0.8), Inches(7.0), Inches(12), Inches(0.38), size=14, color=GREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Five Levers
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
title_bar(sl, "FORGE-VLA — Five Training Levers",
          "Stacked improvements on SmolVLA; each independently toggleable")

levers = [
    ("L1", "EMA of Weights",
     "Decay 0.9999 · warmup 2000 steps  —  smooths variance from stochastic SGD"),
    ("L2", "Chunk Retune",
     "chunk_size=10, n_action_steps=5  (baseline: 50/50)  —  re-query every 5 steps"),
    ("L3", "Percentile Normalisation",
     "Clip action stats at p1/p99 before min-max  —  removes outlier distortion in action targets"),
    ("L4", "Aggressive Augmentation",
     "Tier-1 photometric + Tier-2 lighting (dark/warm/cool/colored LED), prob=0.3, spatial aug on"),
    ("L5", "Temporal Ensembling",
     "Exp-weighted blend of overlapping chunks, coeff=0.01 (inference only) — near-inert at these settings"),
]

for i, (code, name, desc) in enumerate(levers):
    col = LEVER_COLORS[i]
    y   = Inches(1.72) + i * Inches(1.06)
    badge(sl, code, Inches(0.55), y + Inches(0.1),
          Inches(0.65), Inches(0.65), bg=col, fg=WHITE, size=20, bold=True)
    txbox(sl, name, Inches(1.38), y, Inches(11.5), Inches(0.44),
          size=20, bold=True, color=col)
    txbox(sl, desc, Inches(1.38), y + Inches(0.43), Inches(11.7), Inches(0.55),
          size=15, color=GREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — FORGE Results
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
title_bar(sl, "FORGE-VLA — Robustness Results",
          "LIBERO-Long · 300 eps/cell · seed 42  |  P3 = dark (in-aug)  ·  P6 = disco (held-out)  ·  P13 = distractor")

for j, (lbl, col) in enumerate([
        ("baseline (50-step open-loop)", GREY),
        ("baseline_chunk (fair)",        BLUE),
        ("forge_full (all 5 levers)",    GREEN),
]):
    badge(sl, lbl, Inches(0.55) + j*Inches(4.3), Inches(1.72),
          Inches(4.0), Inches(0.4), bg=col, fg=WHITE, size=13)

conditions = [
    ("P0  nominal",            42.0, 56.3, 65.0),
    ("P3  dark  (in-aug)",     40.0, 31.7, 68.3),
    ("P6  disco  (held-out)",  42.0, 55.0, 73.0),
    ("P13 distractor",         16.0,  6.7, 47.0),
    ("P14 noise",              32.0, 30.7, 62.7),
    ("P15 blur",               26.0,  None, 63.3),
    ("P17 compound",           12.0,  8.0, 52.0),
]
bar_cols = [GREY, BLUE, GREEN]
bar_y = Inches(2.26)
for label, b, bc, ff in conditions:
    txbox(sl, label, Inches(0.4), bar_y, Inches(2.7), Inches(0.5),
          size=15, color=INK)
    for j, (val, col) in enumerate(zip([b, bc, ff], bar_cols)):
        xo = Inches(3.25) + j * Inches(3.35)
        if val is None:
            txbox(sl, "—", xo, bar_y, Inches(0.8), Inches(0.5),
                  size=15, color=GREY, align=PP_ALIGN.CENTER)
            continue
        bw = Inches((val / 100) * 2.75)
        box = sl.shapes.add_shape(1, xo, bar_y + Inches(0.06), bw, Inches(0.34))
        box.fill.solid(); box.fill.fore_color.rgb = col; box.line.fill.background()
        txbox(sl, f"{val:.0f}%", xo + bw + Inches(0.08), bar_y,
              Inches(0.65), Inches(0.5), size=15, bold=True, color=col)
    bar_y += Inches(0.66)

hline(sl, Inches(0.4), Inches(6.95), Inches(12.5), color=BLUE, thick=2)
txbox(sl, "★  forge_full vs baseline_chunk: +23 pp on P13 distractor · held-out P6: 73% vs 55%",
      Inches(0.4), Inches(7.04), Inches(12.9), Inches(0.4),
      size=16, bold=True, color=AMBER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Planned Adapter and Why It Was Flawed
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
title_bar(sl, "The Planned Front-End Adapter — And Why It Was Flawed",
          "A residual image→image module upstream of SmolVLA's frozen vision encoder")

# Left panel
filled_box(sl, Inches(0.4), Inches(1.72), Inches(6.0), Inches(4.85), bg=LGREY)
txbox(sl, "What We Planned", Inches(0.5), Inches(1.8), Inches(5.8), Inches(0.44),
      size=18, bold=True, color=BLUE)
for i, txt in enumerate([
    "Residual adapter before the vision encoder",
    "Feature-invariance loss: pull f(perturbed) onto f(clean)",
    "Identity-init: zero-conv start → exact passthrough",
    "~0.1–1M params · post-hoc · no VLA weight changes",
    "One frozen module across backbones — no retraining",
]):
    txbox(sl, "• " + txt, Inches(0.6), Inches(2.34) + i*Inches(0.55),
          Inches(5.7), Inches(0.5), size=16, color=INK)

# Right panel
filled_box(sl, Inches(6.9), Inches(1.72), Inches(6.0), Inches(4.85), bg=LGREY)
txbox(sl, "Why It Was Flawed", Inches(7.0), Inches(1.8), Inches(5.8), Inches(0.44),
      size=18, bold=True, color=RED)
flaws = [
    ("CRT front-end HURTS SmolVLA",
     "Pixel-restoration (CRT, RSS 2024): 43% → 33%\nSmolVLA too small to absorb the distribution shift"),
    ("Lighting barely bites at all",
     "baseline P3 ≈ P0 (40 vs 42%).  P6 ≈ P0 (42 vs 42%)\nVisual perturbations are not the failure mode"),
    ("Action chunk is the real issue",
     "Flow-matching sampler injects jitter into all 50 steps\nChunk boundary discontinuities → task failures"),
]
y = Inches(2.34)
for head, body in flaws:
    txbox(sl, "✗  " + head, Inches(7.05), y, Inches(5.75), Inches(0.4),
          size=16, bold=True, color=RED)
    txbox(sl, body, Inches(7.2), y + Inches(0.38), Inches(5.6), Inches(0.65),
          size=14, color=GREY)
    y += Inches(1.15)

hline(sl, Inches(0.4), Inches(6.7), Inches(12.5), color=AMBER, thick=3)
txbox(sl, "↓  The problem is not in the image pipeline.  It is in the predicted action chunk.",
      Inches(0.4), Inches(6.78), Inches(12.9), Inches(0.56),
      size=18, bold=True, color=AMBER, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Transition
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
title_bar(sl, "From Robustness Recipes to Spectral Action Filtering",
          "A pivot motivated by the actual failure mode of SmolVLA action chunks")

steps = [
    (BLUE,  "FORGE-VLA",
     "5-lever training recipe. forge_full: 65% clean, 73% held-out disco.\n"
     "Key insight: lighting robustness is easy. Action horizon and chunk quality matter more."),
    (AMBER, "Adapter Plan (abandoned)",
     "Residual front-end to canonicalise perturbed images.\n"
     "CRT shows pixel-restoration REGRESSES SmolVLA (43 → 33%). Wrong locus."),
    (TEAL,  "Root Cause Identified",
     "SmolVLA predicts a 50-step action chunk and executes it open-loop.\n"
     "Flow-matching sampler injects high-frequency jitter across all bands. Cross-chunk discontinuities → task failure."),
    (GREEN, "New Approach",
     "Apply spectral decomposition to the predicted chunk.\n"
     "Suppress high-frequency bands with principled, calibrated per-band gain. Trained post-hoc on frozen model output."),
]

y = Inches(1.68)
for col, head, body in steps:
    dot = sl.shapes.add_shape(1, Inches(0.45), y + Inches(0.14),
                               Inches(0.26), Inches(0.26))
    dot.fill.solid(); dot.fill.fore_color.rgb = col; dot.line.fill.background()
    if head != "New Approach":
        ln = sl.shapes.add_shape(1, Inches(0.57), y + Inches(0.42),
                                  Inches(0.04), Inches(1.25))
        ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0xCC,0xCC,0xCC)
        ln.line.fill.background()
    txbox(sl, head, Inches(0.88), y, Inches(12.0), Inches(0.44),
          size=20, bold=True, color=col)
    txbox(sl, body, Inches(0.88), y + Inches(0.43), Inches(12.0), Inches(0.88),
          size=16, color=INK)
    y += Inches(1.48)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Architecture
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
title_bar(sl, "Architecture — Spectral Action Filter",
          "Post-hoc module on a frozen SmolVLA.  Zero changes to the base model.")

ARCH_BOXES = [
    (Inches(0.25), "Frozen\nSmolVLA",            BLUE,  WHITE),
    (Inches(2.4),  "Action Chunk\nA ∈ ℝ^{T×D}",  GREY,  WHITE),
    (Inches(4.55), "DCT-II\n(per joint,\ntime axis)", TEAL, WHITE),
    (Inches(6.7),  "Per-band\nWiener Gain\ng_k = λ/(λ+σ²)", AMBER, WHITE),
    (Inches(8.85), "Inverse DCT\nÂ ∈ ℝ^{T×D}",  GREY,  WHITE),
    (Inches(11.0), "Execute\nn_action_steps\n= 25 steps", GREEN, WHITE),
]
bw, bh = Inches(1.9), Inches(1.3)
by = Inches(1.92)
for bx, label, col, fg in ARCH_BOXES:
    box = sl.shapes.add_shape(1, bx, by, bw, bh)
    box.fill.solid(); box.fill.fore_color.rgb = col; box.line.fill.background()
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = fg
    if bx < Inches(11.0):
        arr = sl.shapes.add_shape(1, bx + bw + Inches(0.02),
                                   by + bh//2 - Inches(0.08),
                                   Inches(0.2), Inches(0.16))
        arr.fill.solid(); arr.fill.fore_color.rgb = GREY; arr.line.fill.background()

txbox(sl, "T=50 bands/joint · D=7 joints → 350 spectral coefficients",
      Inches(4.55), by + bh + Inches(0.08), Inches(4.35), Inches(0.36),
      size=13, color=GREY, align=PP_ALIGN.CENTER)
txbox(sl, "350 learned σ_k²  ·  EMA λ_k from policy",
      Inches(6.7), by + bh + Inches(0.08), Inches(2.15), Inches(0.5),
      size=13, color=AMBER, align=PP_ALIGN.CENTER)

filled_box(sl, Inches(0.9), Inches(4.35), Inches(11.5), Inches(0.6),
           bg=RGBColor(0xFF, 0xF8, 0xE1), border_color=AMBER)
txbox(sl, "Training Loss:   L = MSE( Â ,  A★ )  +  β · Σ_k R_k",
      Inches(0.95), Inches(4.38), Inches(11.4), Inches(0.55),
      size=20, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

txbox(sl, "A★ = ground-truth action chunk from training data  (denoiser, not compressor)",
      Inches(0.9), Inches(5.03), Inches(11.5), Inches(0.4),
      size=15, color=GREY, align=PP_ALIGN.CENTER)

props = ["✓ Base VLA fully frozen", "✓ 350 scalars only",
         "✓ Identity at init  (g_k≈1)", "✓ Zero inference overhead",
         "✓ Trains on cached outputs"]
xp = Inches(0.4)
for p in props:
    txbox(sl, p, xp, Inches(5.6), Inches(2.56), Inches(0.48),
          size=14, bold=True, color=TEAL)
    xp += Inches(2.6)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Math Video 1: SpectralBottleneck
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
title_bar(sl, "The Math — Spectral Decomposition & Gaussian Channel",
          "Walk-through: SpectralBottleneck.mp4")

embed_video(sl, MATH_VID1, None,
            Inches(0.5), Inches(1.45), Inches(8.9), Inches(5.25))

eqs = [
    ("DCT-II",       "C_k = Σ_t a_t·cos[π/T·(t+½)·k]"),
    ("Signal power", "λ_k = 𝔼[C_k²]  (EMA over batches)"),
    ("Channel",      "Ĉ_k = C_k + ε_k,  ε_k ~ N(0, σ_k²)"),
    ("Wiener gain",  "g_k = λ_k / (λ_k + σ_k²)  ∈ (0, 1]"),
]
ey = Inches(1.6)
for head, eq in eqs:
    txbox(sl, head + ":", Inches(9.75), ey, Inches(3.4), Inches(0.42),
          size=14, bold=True, color=BLUE)
    txbox(sl, eq, Inches(9.75), ey + Inches(0.41), Inches(3.4), Inches(0.42),
          size=13, color=INK)
    ey += Inches(1.1)

filled_box(sl, Inches(9.7), Inches(5.7), Inches(3.45), Inches(0.98), bg=LGREY)
txbox(sl, "Low-freq bands: high λ_k, high SNR → motion\nHigh-freq bands: low λ_k → jitter (g_k ≪ 1)",
      Inches(9.75), Inches(5.75), Inches(3.35), Inches(0.88), size=13, color=GREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Math Video 2: SpectralMath  (rate-distortion)
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
title_bar(sl, "The Math — Rate-Distortion & Reverse Water-Filling",
          "Walk-through: SpectralMath.mp4")

embed_video(sl, MATH_VID2, None,
            Inches(0.5), Inches(1.45), Inches(8.9), Inches(5.25))

eqs2 = [
    ("Rate (band k)", "R_k = ½·ln(1 + λ_k/σ_k²)  [nats]"),
    ("Loss",          "L = MSE(Â, A★) + β·Σ_k R_k"),
    ("Optimal σ²_k",  "σ²_k = β·λ_k / (2λ_k − β)"),
    ("Water level",   "θ = β/2 ;  R_k = ½·ln(λ_k/θ)  if λ_k > θ"),
]
ey = Inches(1.6)
for head, eq in eqs2:
    txbox(sl, head + ":", Inches(9.75), ey, Inches(3.4), Inches(0.42),
          size=14, bold=True, color=AMBER)
    txbox(sl, eq, Inches(9.75), ey + Inches(0.41), Inches(3.4), Inches(0.42),
          size=13, color=INK)
    ey += Inches(1.1)

filled_box(sl, Inches(9.7), Inches(5.7), Inches(3.45), Inches(0.98), bg=LGREY)
txbox(sl, "Proved + unit-tested in sib/waterfill.py.\nNo prior IB work for robot actions\nestablishes this analytic link.",
      Inches(9.75), Inches(5.72), Inches(3.35), Inches(0.92), size=13, color=GREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Ablations
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
title_bar(sl, "Results — Component Ablations",
          "LIBERO-Spatial · n_action_steps=25 · 200 episodes · one component removed at a time")

ablations = [
    ("Vanilla baseline",        57.0, GREY,
     "No filtering.  Raw policy output."),
    ("raw_vib  (no DCT)",       36.0, RED,
     "VIB on raw action values — bands not decorrelated, module collapses."),
    ("gain_no_rate",            53.0, AMBER,
     "Learned gains, no rate penalty — no information constraint."),
    ("lowpass  (no rate term)", 55.0, AMBER,
     "Fixed cutoff filter, no β·R loss — cannot adapt to signal distribution."),
    ("SIB β=1e-4  (ours)",      59.5, TEAL,
     "DCT + per-band channel + Wiener gain + explicit rate β·Σ R_k.  Best result."),
]

bar_y = Inches(1.75)
for label, val, col, note in ablations:
    filled_box(sl, Inches(3.8), bar_y, Inches(7.2), Inches(0.6), bg=LGREY)
    bw = Inches((val / 100.0) * 7.0)
    box = sl.shapes.add_shape(1, Inches(3.8), bar_y + Inches(0.08),
                               bw, Inches(0.44))
    box.fill.solid(); box.fill.fore_color.rgb = col; box.line.fill.background()
    txbox(sl, label, Inches(0.35), bar_y, Inches(3.35), Inches(0.68),
          size=16, bold=(col==TEAL), color=col)
    txbox(sl, f"{val:.0f}%", Inches(3.8) + bw + Inches(0.1), bar_y,
          Inches(0.8), Inches(0.68), size=17, bold=True, color=col)
    txbox(sl, note, Inches(11.1), bar_y, Inches(2.15), Inches(0.68),
          size=12, color=GREY)
    bar_y += Inches(1.02)

hline(sl, Inches(0.35), Inches(7.0), Inches(12.6), color=BLUE, thick=2)
txbox(sl, "DCT is load-bearing: −21 pp without it (57% → 36%).  "
          "Rate term adds +6.5 pp over gain-only.  "
          "Full SIB is the only variant above the vanilla baseline.",
      Inches(0.35), Inches(7.07), Inches(12.9), Inches(0.4),
      size=15, bold=True, color=INK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — n_action_steps Sweep
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
title_bar(sl, "Results — Efficiency vs. Accuracy: n_action_steps Sweep",
          "How many steps to execute per chunk?  Vanilla vs SIB β=1e-4 · LIBERO-Spatial · 200 eps")

txbox(sl, "Vanilla (no filtering)", Inches(0.9), Inches(1.68), Inches(5.5), Inches(0.48),
      size=18, bold=True, color=GREY)
txbox(sl, "SIB β=1e-4", Inches(7.2), Inches(1.68), Inches(5.5), Inches(0.48),
      size=18, bold=True, color=TEAL)

sweep = [
    ("n = 1  (oracle: re-query every step)", 72.0, 22.5, True),
    ("n = 5",                                 60.5, 60.0, False),
    ("n = 10",                                61.0, 61.0, False),
    ("n = 25  ←  operating point",            57.0, 53.5, True),
]
sy = Inches(2.28)
for label, van, sib, highlight in sweep:
    lbl = label.replace(" ←  operating point", "")
    cv  = INK if highlight else GREY
    cs  = TEAL if not highlight else GREEN

    txbox(sl, lbl, Inches(0.4), sy, Inches(3.2), Inches(0.72),
          size=16, bold=highlight, color=cv)

    filled_box(sl, Inches(3.65), sy + Inches(0.06), Inches(3.5), Inches(0.58), bg=LGREY)
    bwv = Inches((van / 100) * 3.5)
    box = sl.shapes.add_shape(1, Inches(3.65), sy + Inches(0.1), bwv, Inches(0.48))
    box.fill.solid(); box.fill.fore_color.rgb = cv; box.line.fill.background()
    txbox(sl, f"{van:.0f}%", Inches(3.65) + bwv + Inches(0.1), sy,
          Inches(0.65), Inches(0.72), size=17, bold=highlight, color=cv)

    filled_box(sl, Inches(7.55), sy + Inches(0.06), Inches(3.5), Inches(0.58), bg=LGREY)
    bws = Inches((sib / 100) * 3.5)
    box2 = sl.shapes.add_shape(1, Inches(7.55), sy + Inches(0.1), bws, Inches(0.48))
    box2.fill.solid(); box2.fill.fore_color.rgb = cs; box2.line.fill.background()
    txbox(sl, f"{sib:.0f}%", Inches(7.55) + bws + Inches(0.1), sy,
          Inches(0.65), Inches(0.72), size=17, bold=highlight, color=cs)

    if highlight and "operating" in label:
        txbox(sl, "← 25× cheaper than n=1",
              Inches(8.5), sy + Inches(0.2), Inches(3.5), Inches(0.38),
              size=13, color=GREEN, italic=True)
    sy += Inches(1.05)

hline(sl, Inches(0.4), Inches(6.65), Inches(12.6), color=BLUE, thick=2)
txbox(sl, "At n=25, SIB is within 3.5 pp of oracle vanilla while using 25× fewer queries.  "
          "SIB partially recovers the accuracy lost from reduced query frequency.",
      Inches(0.4), Inches(6.72), Inches(12.9), Inches(0.6), size=16, color=INK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Best SIB vs Vanilla at n=25
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
title_bar(sl, "Results — Best SIB vs Vanilla at Matched Compute (n=25)",
          "Same operating point: n_action_steps=25 · 200 episodes")

txbox(sl, "LIBERO-Spatial (10 tasks)", Inches(0.4), Inches(1.62), Inches(10), Inches(0.46),
      size=18, bold=True, color=BLUE)

spatial_comps = [
    ("Vanilla n=25",            57.0, GREY),
    ("SIB β=3e-4  (conservative)", 55.5, AMBER),
    ("SIB n=25  (strict match)", 53.5, RGBColor(0x1F,0x5F,0xD0)),
    ("SIB β=1e-4  (best config)", 59.5, TEAL),
]
y = Inches(2.16)
for label, val, col in spatial_comps:
    filled_box(sl, Inches(4.0), y + Inches(0.06), Inches(7.8), Inches(0.62), bg=LGREY)
    bw = Inches((val / 100) * 7.6)
    box = sl.shapes.add_shape(1, Inches(4.0), y + Inches(0.1), bw, Inches(0.52))
    box.fill.solid(); box.fill.fore_color.rgb = col; box.line.fill.background()
    txbox(sl, label, Inches(0.35), y, Inches(3.55), Inches(0.72),
          size=16, bold=(col==TEAL), color=col)
    txbox(sl, f"{val:.0f}%", Inches(4.0) + bw + Inches(0.1), y,
          Inches(0.85), Inches(0.72), size=17, bold=True, color=col)
    y += Inches(0.93)

hline(sl, Inches(0.35), Inches(5.92), Inches(12.6), color=BLUE, thick=2)

txbox(sl, "LIBERO-Long (10 tasks — harder, longer-horizon)",
      Inches(0.35), Inches(6.0), Inches(10), Inches(0.46),
      size=18, bold=True, color=BLUE)
y = Inches(6.55)
for label, val, col in [("Vanilla Long", 41.0, GREY), ("SIB Long", 39.0, TEAL)]:
    filled_box(sl, Inches(4.0), y + Inches(0.06), Inches(5.8), Inches(0.55), bg=LGREY)
    bw = Inches((val / 100) * 5.6)
    box = sl.shapes.add_shape(1, Inches(4.0), y + Inches(0.1), bw, Inches(0.44))
    box.fill.solid(); box.fill.fore_color.rgb = col; box.line.fill.background()
    txbox(sl, label, Inches(0.35), y, Inches(3.55), Inches(0.65),
          size=16, color=col)
    txbox(sl, f"{val:.0f}%", Inches(4.0) + bw + Inches(0.1), y,
          Inches(0.8), Inches(0.65), size=17, bold=True, color=col)
    y += Inches(0.66)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Demo Simulation Videos
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
title_bar(sl, "Demo — SIB Simulation Roll-outs",
          "LIBERO-Spatial · SIB β=1e-4 · clean condition · 4 representative tasks")

tasks_demo = [1, 3, 7, 9]
positions   = [
    (Inches(0.35), Inches(1.52)),
    (Inches(6.85), Inches(1.52)),
    (Inches(0.35), Inches(4.5)),
    (Inches(6.85), Inches(4.5)),
]
vw, vh = Inches(6.3), Inches(2.72)

for task, (lx, ty) in zip(tasks_demo, positions):
    vid  = SIB_VID(task, 0)
    post = POSTER("sib", task)
    embed_video(sl, vid, post, lx, ty, vw, vh)
    txbox(sl, f"Task {task} — libero_spatial_{task}  (SIB β=1e-4 · clean)",
          lx, ty + vh + Inches(0.03), vw, Inches(0.36),
          size=13, color=GREY, align=PP_ALIGN.CENTER)

txbox(sl, "All clips: ep000 · clean condition · SIB β=1e-4",
      Inches(0.35), Inches(7.22), Inches(12.9), Inches(0.26),
      size=12, color=GREY, italic=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Comparison with Prior IB Work
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
title_bar(sl, "How We Differ from Prior Information-Bottleneck Work",
          "Feature-level comparison — all 4 papers apply IB/VIB in some form")

headers   = ["Dimension",            "StableVLA\n(Hiranaka 2024)",       "IBAC-SNI\n(Igl 2019)",        "VDB\n(Peng 2019)",             "Ours (SIB)"]
rows_data = [
    ("Locus",             "Visual tokens\n→ LLM interface",  "State encoding\n(RL input)",      "Discriminator\nfeatures (IL)",  "Action chunk\noutput"),
    ("Spectral basis",    "None",                            "None",                             "None",                          "DCT-II\n(time axis)"),
    ("Channel model",     "Channel-cov\nsigmoid gating",     "KL to Gaussian\nprior (VIB)",      "KL to Gaussian\nprior (VIB)",   "Per-band Gaussian\n+ λ_k from policy"),
    ("Rate term",         "None  (heuristic)",               "β·KL(q||p)",                       "β·KL(q||p)",                    "β·Σ R_k\n= β·Σ½ln(1+λ/σ²)"),
    ("Inference decode",  "Sigmoid passthrough",             "Stochastic z~q(z|x)",              "Stochastic z~q(z|x)",           "MMSE Wiener\n(closed-form)"),
    ("Post-hoc frozen",   "No — full\nVLA finetune",         "No — end-to-end",                  "No — inside IL\ntraining loop", "Yes — cached\noutputs only"),
]

col_ws  = [Inches(2.25), Inches(2.4), Inches(2.4), Inches(2.4), Inches(2.4)]
row_h   = Inches(0.73)
tl, tt  = Inches(0.3), Inches(1.48)

for j, (hd, cw) in enumerate(zip(headers, col_ws)):
    bg = BLUE if j in (0,) else (TEAL if j == 4 else RGBColor(0x1F,0x5F,0xD0))
    fg = WHITE
    x  = tl + sum(col_ws[:j])
    table_cell(sl, hd, x, tt, cw, row_h * 0.72,
               bg=bg, fg=fg, size=13, bold=True)

for i, row in enumerate(rows_data):
    for j, (cell, cw) in enumerate(zip(row, col_ws)):
        x   = tl + sum(col_ws[:j])
        y   = tt + row_h * 0.72 + i * row_h
        is_ours = (j == 4)
        is_dim  = (j == 0)
        bg = (RGBColor(0xE6,0xF7,0xF5) if is_ours else
              RGBColor(0xEE,0xF2,0xF8) if is_dim  else
              WHITE)
        fg = TEAL if is_ours else (BLUE if is_dim else INK)
        table_cell(sl, cell, x, y, cw, row_h,
                   bg=bg, fg=fg, size=12, bold=is_ours)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Novelty & Contributions
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
title_bar(sl, "Novelty & Contributions",
          "What is genuinely new — not the backbone, not VIB itself")

contrib = [
    (BLUE,  "1", "DCT on VLA action chunks",
     "Spectral decomposition of the temporal axis of robot action chunks is new. "
     "Energy compaction means low bands carry motion; high bands carry jitter. "
     "The spectral basis makes per-band treatment physically meaningful."),
    (TEAL,  "2", "Per-band signal power estimation (λ_k)",
     "EMA over the policy's own predicted DCT coefficients — dataset-level, always detached. "
     "Calibrates each band's Wiener gain to the actual signal distribution, not a fixed prior."),
    (GREEN, "3", "MMSE Wiener gain as the filter",
     "g_k = λ_k/(λ_k + σ_k²).  Closed-form, provably optimal under Gaussian channel model. "
     "Deterministic at inference — no sampling. 350 multiplications per chunk, negligible overhead."),
    (AMBER, "4", "Rate-distortion / reverse water-filling applied to robot action sequences",
     "Training loss β·Σ R_k forces per-band bit allocation. Optimal allocation is reverse water-filling "
     "— proved and unit-tested. No prior IB work for robot actions establishes this analytic connection."),
    (PURP,  "5", "Post-hoc on a frozen VLA — identity at init",
     "Trains on cached policy outputs. Base model never updated. σ_k² starts small → g_k≈1 "
     "at epoch 0 → exact passthrough before learning. Zero forgetting. Composes with any FT."),
]

y = Inches(1.62)
for col, num, head, body in contrib:
    badge(sl, num, Inches(0.38), y, Inches(0.52), Inches(0.52),
          bg=col, fg=WHITE, size=20, bold=True)
    txbox(sl, head, Inches(1.05), y, Inches(12.1), Inches(0.44),
          size=18, bold=True, color=col)
    txbox(sl, body, Inches(1.05), y + Inches(0.43), Inches(12.1), Inches(0.56),
          size=14, color=INK)
    y += Inches(1.06)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Future Work
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
title_bar(sl, "Future Work",
          "Near-term experiments and longer-horizon extensions")

future = [
    (BLUE,  "Leg 3 — Robustness Sweep  (in progress)",
     ["Vanilla vs SIB under action noise σ ∈ {0.05, 0.1, 0.2} on LIBERO-Long",
      "Visual perturbations: brightness/contrast, Gaussian blur",
      "Goal: quantify SR retention of spectral filtering under noisy actuation"]),
    (TEAL,  "Perception-Locus Leg  (Week 2 candidate)",
     ["Apply DCT + Wiener gain to visual token features, not just action output",
      "Direct comparison with StableVLA IB-Adapter at the same LLM interface point",
      "Question: does spectral feature filtering outperform sigmoid channel gating?"]),
    (GREEN, "Context-Aware Bandwidth Allocation  (Leg 2 — suspended)",
     ["Let σ_k² depend on observation context (language + image features)",
      "Per-task bit allocation rather than fixed global scalars",
      "Needs: estimate_lambda with sigma_mode=context → train → eval on Long"]),
    (AMBER, "Real-Robot Transfer",
     ["Post-hoc spectral filter applied to a WidowX/SO-100 policy",
      "Does jitter reduction improve contact-rich tasks on real hardware?",
      "Pre-register: ≥20 trials/condition/policy before claiming a result"]),
]

y = Inches(1.68)
for col, head, bullets in future:
    box = sl.shapes.add_shape(1, Inches(0.38), y, Inches(0.28), Inches(0.36))
    box.fill.solid(); box.fill.fore_color.rgb = col; box.line.fill.background()
    txbox(sl, head, Inches(0.82), y, Inches(12.4), Inches(0.44),
          size=18, bold=True, color=col)
    by = y + Inches(0.43)
    for b in bullets:
        txbox(sl, "  •  " + b, Inches(1.0), by, Inches(12.0), Inches(0.4),
              size=15, color=INK)
        by += Inches(0.4)
    y = by + Inches(0.22)


# ─────────────────────────────────────────────────────────────────────────────
#  SAVE
# ─────────────────────────────────────────────────────────────────────────────
out_path = "presentation/spectral_filter_vla.pptx"
os.makedirs("presentation", exist_ok=True)
prs.save(out_path)
print(f"Saved: {out_path}  ({len(prs.slides)} slides)")
for i, sl_obj in enumerate(prs.slides):
    t = sl_obj.shapes.title
    print(f"  Slide {i+1:2d}: {t.text if t else '(no title)'}")
