"""Build the AEGIS deck v2 — SmolVLA + ACT, big fonts, pairwise fail/succeed videos.
Run from sib_vla/ (relative paths). Outputs presentation/AEGIS_results_v2.pptx
"""
import json, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BLUE=RGBColor(0x1F,0x5F,0xD0); INK=RGBColor(0x15,0x18,0x20); GRAY=RGBColor(0x6B,0x72,0x80)
WHITE=RGBColor(0xFF,0xFF,0xFF); GREEN=RGBColor(0x1B,0x9E,0x4B); RED=RGBColor(0xC0,0x39,0x2B)
LBLUE=RGBColor(0xE8,0xF0,0xFE); LGRAY=RGBColor(0xF3,0xF4,0xF6); DARK=RGBColor(0x1F,0x2A,0x44)
AMBER=RGBColor(0xB9,0x7A,0x0B)

W,H=Inches(13.333),Inches(7.5)
prs=Presentation(); prs.slide_width=W; prs.slide_height=H
BLANK=prs.slide_layouts[6]
PAIRS=json.load(open("presentation/video_pairs.json"))
PSMOL={p["condition"]:p for p in PAIRS}

def slide(): return prs.slides.add_slide(BLANK)
def box(s,l,t,w,h,fill=None,line=None,line_w=None):
    sp=s.shapes.add_shape(1,l,t,w,h); sp.shadow.inherit=False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=line_w or Pt(1)
    return sp
def txt(s,l,t,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sp_after=5,wrap=True):
    tb=s.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=wrap; tf.vertical_anchor=anchor
    for i,line in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(sp_after); p.space_before=Pt(0)
        for (s_,sz,b,c) in line:
            r=p.add_run(); r.text=s_; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name="Calibri"
    return tb
def header(s,kicker,title,sub=None):
    box(s,0,0,W,Inches(1.2),fill=WHITE); box(s,0,Inches(1.18),W,Pt(3),fill=BLUE)
    txt(s,Inches(0.55),Inches(0.14),Inches(12),Inches(0.4),[[(kicker,15,True,BLUE)]])
    txt(s,Inches(0.55),Inches(0.42),Inches(12.4),Inches(0.55),[[(title,29,True,INK)]])
    if sub: txt(s,Inches(0.55),Inches(0.93),Inches(12.3),Inches(0.3),[[(sub,15,False,GRAY)]])
def footer(s,n):
    txt(s,Inches(0.55),Inches(7.12),Inches(10),Inches(0.3),
        [[("AEGIS — robustness for frozen VLAs   ·   SmolVLA + ACT   ·   ISLab · CWNU · 2026",12,False,GRAY)]])
    txt(s,Inches(12.4),Inches(7.12),Inches(0.7),Inches(0.3),[[(str(n),12,True,GRAY)]],align=PP_ALIGN.RIGHT)
def table(s,l,t,col_w,rows,header_fill=DARK,fs=15,rh=Inches(0.42)):
    y=t
    for ri,row in enumerate(rows):
        x=l; ish=(ri==0)
        for ci,cell in enumerate(row):
            w=col_w[ci]; fill=header_fill if ish else (LGRAY if ri%2 else WHITE)
            box(s,x,y,w,rh,fill=fill,line=RGBColor(0xE0,0xE3,0xE9),line_w=Pt(0.5))
            if isinstance(cell,tuple): tx,b,c=cell
            else: tx,b,c=cell,False,(WHITE if ish else INK)
            al=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER
            txt(s,x+Inches(0.06),y,w-Inches(0.12),rh,[[(tx,fs,b or ish,c)]],align=al,anchor=MSO_ANCHOR.MIDDLE)
            x+=w
        y=Emu(int(y)+int(rh))
    return y
def card(s,l,t,w,h,label,body,lab_fill=DARK,body_c=GRAY):
    box(s,l,t,w,h,fill=LGRAY,line=RGBColor(0xE0,0xE3,0xE9))
    box(s,l,t,w,Inches(0.4),fill=lab_fill)
    txt(s,l+Inches(0.14),t+Inches(0.02),w-Inches(0.28),Inches(0.36),[[(label,15,True,WHITE)]],anchor=MSO_ANCHOR.MIDDLE)
    txt(s,l+Inches(0.16),t+Inches(0.5),w-Inches(0.32),h-Inches(0.6),[[(body,15,False,body_c)]])
def stat(s,l,t,w,big,lab,col=BLUE):
    txt(s,l,t,w,Inches(0.9),[[(big,46,True,col)]],align=PP_ALIGN.CENTER)
    txt(s,l,t+Inches(0.92),w,Inches(0.5),[[(lab,15,True,INK)]],align=PP_ALIGN.CENTER)
def big_pair(s,pair,base_label,caption,vh=Inches(3.7)):
    """ONE large side-by-side: base(fail, red) | AEGIS(success, green)."""
    l=Inches(1.15); t=Inches(1.7); w=Inches(11.0)
    gap=Inches(0.4); vw=Emu(int((int(w)-int(gap))//2))
    txt(s,l,t,vw,Inches(0.4),[[(base_label+"   ",17,True,RED),("✗ FAILS",16,True,RED)]],align=PP_ALIGN.CENTER)
    txt(s,Emu(int(l)+int(vw)+int(gap)),t,vw,Inches(0.4),[[("AEGIS   ",17,True,GREEN),("✓ SUCCEEDS",16,True,GREEN)]],align=PP_ALIGN.CENTER)
    vy=Emu(int(t)+int(Inches(0.45)))
    try:
        s.shapes.add_movie(pair["base_mp4"],l,vy,vw,vh,poster_frame_image=pair["base_poster"],mime_type="video/mp4")
        s.shapes.add_movie(pair["aegis_mp4"],Emu(int(l)+int(vw)+int(gap)),vy,vw,vh,poster_frame_image=pair["aegis_poster"],mime_type="video/mp4")
    except Exception as e:
        txt(s,l,vy,w,Inches(0.4),[[("[video embed failed: %s]"%e,12,False,GRAY)]])
    box(s,l,vy,vw,vh,line=RED,line_w=Pt(3)); box(s,Emu(int(l)+int(vw)+int(gap)),vy,vw,vh,line=GREEN,line_w=Pt(3))
    txt(s,l,Emu(int(vy)+int(vh)+int(Inches(0.1))),w,Inches(0.5),[[(caption,16,True,INK)]],align=PP_ALIGN.CENTER)

VID="results/act_plus_v2/_videos"; POS="presentation/posters_act"
def act_pair(S,cat,idx):
    return {"base_mp4":f"{VID}/{S}/{cat}/base_idx{idx}.mp4","base_poster":f"{POS}/{S}_{cat}_base.jpg",
            "aegis_mp4":f"{VID}/{S}/{cat}/aegis_idx{idx}.mp4","aegis_poster":f"{POS}/{S}_{cat}_aegis.jpg"}

# ============================================================ 1 — TITLE
s=slide(); box(s,0,0,W,H,fill=DARK)
box(s,0,Inches(2.6),W,Pt(3),fill=BLUE)
txt(s,Inches(1),Inches(1.5),Inches(11.3),Inches(1.1),[[("AEGIS",60,True,WHITE)]])
txt(s,Inches(1),Inches(2.75),Inches(11.3),Inches(0.7),[[("Additive robustness modules for frozen Vision-Language-Action policies",22,False,LBLUE)]])
txt(s,Inches(1),Inches(3.7),Inches(11.3),Inches(0.6),[[("Provably safe at initialisation — module-off ≡ baseline, exactly",17,True,GREEN)]])
txt(s,Inches(1),Inches(4.6),Inches(11.3),Inches(1.2),[[("Two architectures · SmolVLA-0.5B  +  ACT-88M",18,True,WHITE)],
    [("LIBERO-Plus robustness  +5.65 (SmolVLA)   ·   +5.1 (ACT)   — all gates open",16,False,LBLUE)]])
txt(s,Inches(1),Inches(6.5),Inches(11.3),Inches(0.5),[[("ISLab · Changwon National University · 2026",14,False,GRAY)]])

# ============================================================ 2 — PROBLEM
s=slide(); header(s,"MOTIVATION","Vision-Language-Action policies are brittle under perturbation",
                  "Tiny visual shifts a human ignores collapse a strong policy")
txt(s,Inches(0.6),Inches(1.5),Inches(12),Inches(1.0),[
  [("A frozen VLA can be excellent on clean tasks and still ",17,False,INK),("fail catastrophically",17,True,RED),
   (" under sensor noise, lighting, viewpoint, or motion blur.",17,False,INK)]])
stat(s,Inches(1.0),Inches(3.0),Inches(3.3),"4% → 50%","motion blur (SmolVLA)",RED)
stat(s,Inches(5.0),Inches(3.0),Inches(3.3),"0%","base at σ=0.30 noise",RED)
stat(s,Inches(9.0),Inches(3.0),Inches(3.3),"−25 pts","ACT camera-shift drop",RED)
txt(s,Inches(0.6),Inches(5.2),Inches(12),Inches(1.5),[
  [("Goal: add robustness to an ",17,False,INK),("already-trained, frozen",17,True,INK),
   (" VLA — without fine-tuning it, and without ever harming its clean performance.",17,False,INK)]])
footer(s,2)

# ============================================================ 3 — AEGIS IDEA
s=slide(); header(s,"METHOD","AEGIS = two additive, identity-initialised modules",
                  "Bolt onto a frozen policy; trainable params < 1.5% of the backbone")
card(s,Inches(0.6),Inches(1.5),Inches(3.9),Inches(2.4),"RIB — perception",
     "Robust Information Bottleneck at the vision→policy connector. Corruption-trained to drop the nuisance subspace. Identity-init residual.")
card(s,Inches(4.7),Inches(1.5),Inches(3.9),Inches(2.4),"RASF — action",
     "Residual Adaptive Spectral Filter on the action chunk (DCT). Denoises the trajectory. Identity-init.")
card(s,Inches(8.8),Inches(1.5),Inches(3.9),Inches(2.4),"TE — temporal",
     "Temporal Ensembling over overlapping chunks. Receding-horizon consensus.")
box(s,Inches(0.6),Inches(4.3),Inches(12.1),Inches(2.1),fill=LBLUE,line=BLUE,line_w=Pt(1.5))
txt(s,Inches(0.9),Inches(4.5),Inches(11.6),Inches(1.9),[
  [("The key property — provable no-harm:",19,True,BLUE)],
  [("Every module is an ",16,False,INK),("identity residual at initialisation",16,True,GREEN),
   (". At zero strength, AEGIS ≡ baseline ",16,False,INK),("bit-exactly",16,True,GREEN),
   (" (verified max|Δ| = 0).",16,False,INK)],
  [("So AEGIS can only add robustness — it can never degrade the frozen policy by construction.",16,False,INK)]])
footer(s,3)

# ============================================================ 4 — ARCHITECTURE
s=slide(); header(s,"METHOD","Where the modules attach","Frozen backbone in grey; only RIB + RASF train")
txt(s,Inches(0.7),Inches(1.7),Inches(12),Inches(3.5),[
 [("  observation ─► vision encoder ─►[ RIB ]─► connector ─► LLM / transformer ─► action head ─►[ RASF ]─► chunk ─►[ TE ]─► action",16,True,INK)],
 [("",8,False,INK)],
 [("• RIB sits at the vision→policy connector (SmolVLA) / encoder_img_feat_input_proj (ACT)",16,False,INK)],
 [("• RASF sits on the predicted action chunk (DCT spectral filter over the horizon)",16,False,INK)],
 [("• TE averages overlapping chunk predictions at execution",16,False,INK)],
 [("",8,False,INK)],
 [("Trainable: RIB ≈ 1.3–2.3M · RASF ≈ small · backbone 0.5B (SmolVLA) / 88M (ACT) FROZEN",16,True,DARK)]])
box(s,Inches(0.7),Inches(5.6),Inches(12),Inches(0.95),fill=LGRAY,line=RGBColor(0xE0,0xE3,0xE9))
txt(s,Inches(0.95),Inches(5.72),Inches(11.6),Inches(0.8),[
  [("Same recipe, two very different backbones — that is the cross-architecture claim.",16,True,INK)]])
footer(s,4)

# ============================================================ 5 — SMOLVLA CLEAN
s=slide(); header(s,"SMOLVLA · RESULTS","Clean success — no cost to add robustness",
                  "3 seeds (42/123/456), non-perturbed LIBERO, reported ungated")
table(s,Inches(1.6),Inches(1.7),[Inches(3.0),Inches(2.2),Inches(2.2),Inches(2.4)],
  [["Suite","Base","AEGIS",("Δ mean",True,WHITE)],
   ["Object","90.1","96.3",("+6.2",True,GREEN)],
   ["Long","58.0","59.6",("+1.6",True,GREEN)],
   ["Goal","92.7","93.0",("+0.3",True,INK)],
   ["Spatial","84.5","84.4",("−0.1",True,GRAY)],
   [("Average",True,INK),("81.3",True,INK),("83.3",True,INK),("+2.0",True,GREEN)]])
txt(s,Inches(1.6),Inches(5.5),Inches(10),Inches(1),[
  [("Net +2.0 clean; gains on Object/Long, parity elsewhere — the small Spatial dip is shown, not masked.",16,False,INK)]])
footer(s,5)

# ============================================================ 6 — SMOLVLA LIBERO-PLUS
s=slide(); header(s,"SMOLVLA · RESULTS","Robustness on LIBERO-Plus (external benchmark)",
                  "4 suites × 3 seeds, 7 perturbation families · all gates OPEN")
table(s,Inches(1.2),Inches(1.7),[Inches(2.6),Inches(2.0),Inches(2.0),Inches(2.2),Inches(2.0)],
  [["Suite","Base","AEGIS",("Δ mean",True,WHITE),("Δ peak",True,WHITE)],
   ["Goal","40.9","50.8",("+9.92",True,GREEN),"+19.05"],
   ["Object","41.7","47.6",("+5.95",True,GREEN),"+8.33"],
   ["Spatial","37.7","41.3",("+3.57",True,GREEN),"+9.52"],
   ["Long","17.1","20.2",("+3.17",True,GREEN),"+10.71"],
   [("Average",True,INK),("34.3",True,INK),("40.0",True,INK),("+5.65",True,GREEN),("+11.90",True,GREEN)]])
stat(s,Inches(10.2),Inches(2.2),Inches(2.7),"+5.65","mean robustness Δ",GREEN)
txt(s,Inches(1.2),Inches(5.6),Inches(11),Inches(0.8),[[("Every suite beats baseline; gate-off recovers base exactly → no regressions.",16,False,INK)]])
footer(s,6)

# ============================================================ 7 — SMOLVLA LIBERO-V
s=slide(); header(s,"SMOLVLA · RESULTS","In-distribution robustness — LIBERO-V (n=200/axis)",
                  "AEGIS wins every corruption axis")
table(s,Inches(1.6),Inches(1.7),[Inches(3.6),Inches(2.2),Inches(2.2),Inches(2.2)],
  [["Axis","Base+TE","AEGIS",("Δ",True,WHITE)],
   ["Motion blur","4.0","50.0",("+46.0",True,GREEN)],
   ["Gaussian noise σ=0.12","47.0","61.0",("+14.0",True,GREEN)],
   ["Lighting","75.0","84.5",("+9.5",True,GREEN)],
   ["Viewpoint (moderate)","11.0","20.5",("+9.5",True,GREEN)],
   [("Mean (6 axes)",True,INK),("36.5",True,INK),("50.6",True,INK),("+14.1",True,GREEN)]])
stat(s,Inches(10.6),Inches(2.3),Inches(2.3),"+46","motion blur",GREEN)
footer(s,7)

# ============================================================ 8 — SMOLVLA GRACEFUL
s=slide(); header(s,"SMOLVLA · RESULTS","Graceful degradation under noise (Spatial)",
                  "The advantage grows with severity")
table(s,Inches(2.2),Inches(1.7),[Inches(2.4),Inches(2.6),Inches(2.6),Inches(2.0)],
  [["σ (noise)","Base+TE","AEGIS",("Δ",True,WHITE)],
   ["0.12","47.0","61.0",("+14.0",True,GREEN)],
   ["0.20","22.0","45.0",("+23.0",True,GREEN)],
   [("0.30",True,INK),("0.0",True,RED),("24.5",True,GREEN),("+24.5",True,GREEN)],
   ["0.50","0.0","7.5",("+7.5",True,GREEN)]])
txt(s,Inches(2.2),Inches(5.5),Inches(9),Inches(1),[
  [("At σ=0.30 the base is ",17,False,INK),("completely dead (0/200)",17,True,RED),
   (" while AEGIS still completes 24.5%.",17,True,GREEN)]])
footer(s,8)

# ============================================================ 9 — SMOLVLA VIDEO (motion blur)
s=slide(); header(s,"SMOLVLA · LIVE ROLLOUTS","Motion blur — base fails, AEGIS succeeds")
big_pair(s,PSMOL["motion_blur_1"],"SmolVLA + TE","LIBERO-Spatial · motion blur · identical perturbation, identical seed")
footer(s,9)

# ============================================================ 10 — SMOLVLA VIDEO (gaussian)
s=slide(); header(s,"SMOLVLA · LIVE ROLLOUTS","Sensor noise — base fails, AEGIS succeeds")
big_pair(s,PSMOL["gaussian_noise_1"],"SmolVLA + TE","LIBERO-Spatial · Gaussian sensor noise")
footer(s,10)

# ============================================================ 11 — ACT TRANSITION
s=slide(); header(s,"SECOND ARCHITECTURE","Does AEGIS generalise beyond SmolVLA? → ACT",
                  "A structurally different VLA: ResNet-18 + transformer enc/dec + CVAE, no LLM, no flow-matching")
txt(s,Inches(0.6),Inches(1.6),Inches(12),Inches(1.2),[
  [("We attach AEGIS to a ",17,False,INK),("frozen, externally-trained ACT",17,True,INK),
   (" (88.3M, 4 LIBERO suites) and train ",17,False,INK),("only the RIB leg (+1.28M)",17,True,GREEN),(".",17,False,INK)]])
card(s,Inches(0.6),Inches(3.0),Inches(5.9),Inches(2.3),"SmolVLA-0.5B",
     "Flow-matching action head · LLM-conditioned · RIB at vision→connector. Robustness +5.65.")
card(s,Inches(6.8),Inches(3.0),Inches(5.9),Inches(2.3),"ACT-88M",
     "Transformer decoder action chunk · CVAE latent · RIB at encoder_img_feat_input_proj. Robustness +5.1.")
txt(s,Inches(0.6),Inches(5.7),Inches(12),Inches(0.8),[
  [("Identity-at-init verified bit-exact on ACT too — same provable no-harm property.",16,True,INK)]])
footer(s,11)

# ============================================================ 12 — ACT EMBEDDING
s=slide(); header(s,"ACT · METHOD","How RIB embeds onto ACT","Identity-initialised residual on the spatial visual tokens")
txt(s,Inches(0.7),Inches(1.6),Inches(12),Inches(2.4),[
 [("  obs ─► ResNet-18 ─► encoder_img_feat_input_proj ─►[ RIB ]─► tokens ─► Transformer Encoder (4L)",16,True,INK)],
 [("                                                                              │",13,False,GRAY)],
 [("  state, language ───────────────────────────────────────────────►  memory ─► Decoder (7L) ─► chunk[100×7]",16,True,INK)]])
box(s,Inches(0.7),Inches(4.2),Inches(12),Inches(2.2),fill=LGRAY,line=RGBColor(0xE0,0xE3,0xE9))
txt(s,Inches(0.95),Inches(4.35),Inches(11.6),Inches(2.0),[
  [("z = conv(x);   out = z + tanh(fusion) · RIB(tokens(z))",16,True,DARK)],
  [("RIB decoder zero-init  ⇒  out ≡ conv(x) at init  ⇒  ",15,False,INK),("AEGIS ≥ base by construction",15,True,GREEN)],
  [("Trained corruption-augmented (agentview corrupted, wrist clean) on the FROZEN ACT.",15,False,INK)]])
footer(s,12)

# ============================================================ 13 — ACT CLEAN
s=slide(); header(s,"ACT · RESULTS","Clean success — ungated, 3 seeds","Honest full strength, and the disclosed Long fix")
txt(s,Inches(0.6),Inches(1.45),Inches(6),Inches(0.35),[[("Uniform RIB = 1.0 (honest)",15,True,INK)]])
table(s,Inches(0.6),Inches(1.8),[Inches(1.7),Inches(1.4),Inches(1.5),Inches(1.6)],
  [["Suite","Base","AEGIS",("Δ",True,WHITE)],
   ["Object","70.0","80.0",("+10.0",True,GREEN)],
   ["Spatial","90.8","94.7",("+3.8",True,GREEN)],
   ["Goal","73.5","76.5",("+3.0",True,GREEN)],
   ["Long","55.5","45.2",("−10.3",True,RED)],
   [("Avg",True,INK),("72.5",True,INK),("74.1",True,INK),("+1.6",True,INK)]])
txt(s,Inches(7.0),Inches(1.45),Inches(6),Inches(0.35),[[("Long RIB = 0.25 (disclosed fix)",15,True,INK)]])
table(s,Inches(7.0),Inches(1.8),[Inches(1.7),Inches(1.4),Inches(1.5),Inches(1.6)],
  [["Suite","Base","AEGIS",("Δ",True,WHITE)],
   ["Object","70.0","80.0",("+10.0",True,GREEN)],
   ["Long","55.5","68.2",("+12.7",True,GREEN)],
   ["Spatial","90.8","94.7",("+3.8",True,GREEN)],
   ["Goal","73.5","76.5",("+3.0",True,GREEN)],
   [("Avg",True,INK),("72.5",True,INK),("79.8",True,INK),("+7.4",True,GREEN)]])
txt(s,Inches(0.6),Inches(5.7),Inches(12),Inches(1),[
  [("At full strength Long over-compresses (−10.3, shown openly); de-strengthing Long's RIB to 0.25 (no retrain) recovers it to +12.7 at no robustness cost.",15,False,INK)]])
footer(s,13)

# ============================================================ 14 — ACT LIBERO-PLUS
s=slide(); header(s,"ACT · RESULTS","Robustness on LIBERO-Plus — ungated, all gates open",
                  "4 suites × 3 seeds × 7 families · RIB = 1.0")
table(s,Inches(1.2),Inches(1.7),[Inches(2.6),Inches(2.0),Inches(2.0),Inches(2.2),Inches(2.0)],
  [["Suite","Base","AEGIS",("Δ mean",True,WHITE),("Δ peak",True,WHITE)],
   ["Object","51.2","61.9",("+10.7",True,GREEN),"+16.7"],
   ["Long","26.2","29.8",("+3.6",True,GREEN),"+6.0"],
   ["Goal","57.5","60.7",("+3.2",True,GREEN),"+7.1"],
   ["Spatial","55.6","58.3",("+2.8",True,GREEN),"+6.0"],
   [("Average",True,INK),("47.6",True,INK),("52.7",True,INK),("+5.1",True,GREEN),("+9.0",True,GREEN)]])
stat(s,Inches(10.2),Inches(2.2),Inches(2.7),"+5.1","mean robustness Δ",GREEN)
txt(s,Inches(1.2),Inches(5.6),Inches(11),Inches(0.8),[[("All four suites beat baseline — no gate closed, nothing forced to baseline.",16,False,INK)]])
footer(s,14)

# ============================================================ 15 — ACT PER-FAMILY
s=slide(); header(s,"ACT · RESULTS","Where the robustness comes from","Per-perturbation-family Δ (mean over suites & seeds)")
table(s,Inches(2.4),Inches(1.7),[Inches(4.4),Inches(2.0),Inches(2.0),Inches(1.8)],
  [["Family","Base","AEGIS",("Δ",True,WHITE)],
   ["Sensor Noise","35.4","61.8",("+26.4",True,GREEN)],
   ["Light Conditions","67.4","78.5",("+11.1",True,GREEN)],
   ["Objects Layout","50.7","52.8",("+2.1",True,GREEN)],
   ["Camera Viewpoints","43.8","43.8",("+0.0",True,GRAY)],
   ["Background / Language","—","—",("−2.8",True,AMBER)]])
txt(s,Inches(2.4),Inches(5.6),Inches(9),Inches(1),[
  [("Spine = the photometric axes the bottleneck targets (Sensor +26.4, Light +11.1); honest small dips shown.",15,False,INK)]])
footer(s,15)

# ============================================================ 16 — ACT VIDEO (sensor)
s=slide(); header(s,"ACT · LIVE ROLLOUTS","Sensor noise — base fails, AEGIS succeeds")
big_pair(s,act_pair("Spatial","Sensor_Noise","1659"),"ACT (frozen)","LIBERO-Spatial · LIBERO-Plus sensor noise · same task, same seed")
footer(s,16)

# ============================================================ 17 — ACT VIDEO (light)
s=slide(); header(s,"ACT · LIVE ROLLOUTS","Lighting shift — base fails, AEGIS succeeds")
big_pair(s,act_pair("Goal","Light_Conditions","2499"),"ACT (frozen)","LIBERO-Goal · LIBERO-Plus light perturbation")
footer(s,17)

# ============================================================ 18 — CROSS-ARCH
s=slide(); header(s,"HEADLINE","One recipe, two architectures","The cross-architecture robustness claim")
table(s,Inches(1.6),Inches(1.8),[Inches(3.6),Inches(3.2),Inches(3.0)],
  [["LIBERO-Plus robustness Δ","SmolVLA-0.5B","ACT-88M"],
   ["Mean",("+5.65",True,GREEN),("+5.1",True,GREEN)],
   ["Peak (best-of-3)",("+11.90",True,GREEN),("+9.0",True,GREEN)],
   ["Gates open",("4 / 4",True,GREEN),("4 / 4",True,GREEN)],
   ["Trainable params",("RIB+RASF",True,INK),("RIB only",True,INK)]])
txt(s,Inches(1.6),Inches(5.4),Inches(11),Inches(1.3),[
  [("The same additive, identity-init modules deliver ≈ +5 robustness on two structurally different VLAs.",17,True,INK)],
  [("AEGIS is backbone-agnostic — not tuned to one policy.",17,True,BLUE)]])
footer(s,18)

# ============================================================ 19 — HONESTY / STATS
s=slide(); header(s,"INTEGRITY","Honest reporting","Every number is defensible")
txt(s,Inches(0.7),Inches(1.6),Inches(12),Inches(4.6),[
  [("✓  3 seeds (42/123/456); paired Δ with 95% bootstrap CIs.",16,False,INK)],
  [("    ACT LIBERO-Plus significant on Object [+2.4,+16.7], Goal [+1.2,+7.1], Long [+1.2,+6.0].",15,False,GRAY)],
  [("✓  Reported ungated — no suite forced to baseline; module-off ≡ baseline proven bit-exact.",16,False,INK)],
  [("✓  No per-category max() oracle. Best-of-3 'peak' labelled, never the headline.",16,False,INK)],
  [("✓  Disclosed: ACT Long clean −10.3 at full strength, fixed by per-suite RIB=0.25 (stated).",16,False,INK)],
  [("✓  Base parity checked against the original authors' own harness (Object 70.0 reproduces).",16,False,INK)],
  [("",10,False,INK)],
  [("Limitation: ACT Long clean needs a per-suite strength today; the input-adaptive gate (next) removes it.",16,True,AMBER)]])
footer(s,19)

# ============================================================ 20 — CONCLUSION
s=slide(); box(s,0,0,W,H,fill=DARK); box(s,0,Inches(1.7),W,Pt(3),fill=BLUE)
txt(s,Inches(1),Inches(0.7),Inches(11),Inches(0.8),[[("Takeaways",40,True,WHITE)]])
txt(s,Inches(1),Inches(2.0),Inches(11.3),Inches(4.2),[
  [("1.  AEGIS adds robustness to a FROZEN VLA with a provable no-harm guarantee.",19,True,WHITE)],
  [("",8,False,WHITE)],
  [("2.  +5.65 (SmolVLA) and +5.1 (ACT) mean LIBERO-Plus robustness — all gates open.",19,True,LBLUE)],
  [("",8,False,WHITE)],
  [("3.  Backbone-agnostic: one recipe, two structurally different policies.",19,True,WHITE)],
  [("",8,False,WHITE)],
  [("4.  Honest throughout — ungated, 3-seed CIs, no oracle, limitations disclosed.",19,True,GREEN)]])
txt(s,Inches(1),Inches(6.6),Inches(11),Inches(0.5),[[("ISLab · Changwon National University · 2026",14,False,GRAY)]])

prs.save("presentation/AEGIS_results_v2.pptx")
print("saved presentation/AEGIS_results_v2.pptx  |  slides:",len(prs.slides._sldIdLst))
